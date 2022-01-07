# -*- coding: utf-8 -*-
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.support.ui import Select
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import time
import datetime
from pptx import Presentation
from pptx.util import Cm,Pt
from pptx.parts.slide import Slide
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from PIL import Image
import urllib.request
import requests
import json
from xml.etree import ElementTree as ET
#--------------------------------1-------------------------------#

def reverseMax(a):
  reverseMax = len(a) - a[::-1].index(max(a)) - 1
  return reverseMax
def NulltoZero(a):
  rep = [ 0.0  if x == None else x for x in a]
  return rep
path = str(os.getcwd())
clpfigPath = path + '\\clipIMG\\'
chrome_options = Options()
chrome_options.add_argument('--headless')
#chrome_options.add_argument('--disable-gpu')
browser = webdriver.Chrome(chrome_options=chrome_options)
#browser = webdriver.Chrome()
browser.set_window_size(1400,1000)
wait = WebDriverWait(browser, 10)
prs = Presentation('wrasbsample.pptx')
#-----修改文字----
#-----封面標題----
slide = prs.slides[0]
body_shape = slide.shapes
body_shape[0].text_frame.text = str(int(time.strftime('%Y'))-1911) +'年'+time.strftime('%m')+'月'+time.strftime('%d')+'日 '+time.strftime('%H')+':00\n 水情研判'
body_shape[0].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
body_shape[0].text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
#-----封面標題----
##-----第2頁文字----
slide = prs.slides[1]
body_shape = slide.shapes
body_shape[2].text_frame.paragraphs[1].text = time.strftime('%m')+'/'+time.strftime('%d')+' '+time.strftime('%H') \
+'時 QPESUMS+WRF模式預估未來72小時轄區水庫集水區預報降雨分析 (~'+time.strftime('%m')+'/'+str(int(time.strftime('%d'))+3)+' '+str(int(time.strftime('%H'))-1)+'時):'
body_shape[2].text_frame.paragraphs[1].font.size = Pt(18)
body_shape[2].text_frame.paragraphs[1].font.name = '微軟正黑體'
datas = []
requestURL = 'http://localhost/wsReservoirInfo'
xml = urllib.request.urlopen(requestURL)
root = ET.fromstring(xml.read())
for data in root[1][0].findall('Table'):
  res_name = data.find('res_name').text
  if res_name in ('曾文水庫' , '牡丹水庫' , '阿公店水庫'):
    name = res_name
    level = data.find('WaterLine').text
    Capacity = data.find('Capacity').text
    CapacityRate = data.find('CapacityRate').text
    data_time = data.find('DATE').text
    res_time = datetime.datetime.strptime(data_time,'%Y-%m-%dT%H:00:00+08:00').strftime('%m/%d %H')
    x = (name,level,Capacity,CapacityRate,res_time)
    datas.append(x)
  else:
    continue
print (datas)
for x in datas:
  if x[0] == '曾文水庫':
    body_shape[2].text_frame.paragraphs[2].text = '曾文水庫('+x[4]+'時，水位'+x[1]+'公尺   蓄水量'+x[2]+'萬噸   蓄水率'+x[3]+'%):'
    body_shape[2].text_frame.paragraphs[2].font.size = Pt(16)
    body_shape[2].text_frame.paragraphs[2].runs[0].font.color.rgb = RGBColor(0, 0, 255)
  elif x[0] == '牡丹水庫':
    body_shape[2].text_frame.paragraphs[4].text = '牡丹水庫('+x[4]+'時，水位'+x[1]+'公尺   蓄水量'+x[2]+'萬噸   蓄水率'+x[3]+'%):'
    body_shape[2].text_frame.paragraphs[4].font.size = Pt(16)
    body_shape[2].text_frame.paragraphs[4].font.color.rgb = RGBColor(0, 0, 255)
  elif x[0] == '阿公店水庫':
    body_shape[2].text_frame.paragraphs[6].text = '阿公店水庫('+x[4]+'時，水位'+x[1]+'公尺   蓄水量'+x[2]+'萬噸   蓄水率'+x[3]+'%):'
    body_shape[2].text_frame.paragraphs[6].font.size = Pt(16)
    body_shape[2].text_frame.paragraphs[6].font.color.rgb = RGBColor(0, 0, 255)

requestURL = 'http://localhost/QPESUMSWRF/'
print (requestURL)
js = urllib.request.urlopen(requestURL)
j_obj = json.loads(js.read())
for key in j_obj['values']:
  if key['code'] == 'Res0016':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    ZW72DATA = [name,sum(value),peak,peaktime]
  if key['code'] == 'Res0022':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    MD72DATA = [name,sum(value),peak,peaktime]
  elif key['code'] == 'Res0020':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    AG72DATA = [name,sum(value),peak,peaktime]
  elif key['code'] == '1730w002':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    ZS72DATA = [name,sum(value),peak,peaktime]
  elif key['code'] == 'weir00002':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    GP72DATA = [name,sum(value),peak,peaktime]
#print (ZW72DATA,MD72DATA,AG72DATA,ZS72DATA,GP72DATA)
body_shape[2].text_frame.paragraphs[3].text = '時雨量最高'+str(ZW72DATA[2])+'mm('+ZW72DATA[3]+'時)，累積降雨量'+str(ZW72DATA[1])+'mm'
body_shape[2].text_frame.paragraphs[3].font.size = Pt(16)
body_shape[2].text_frame.paragraphs[3].font.color.rgb = RGBColor(255, 0, 0)
body_shape[2].text_frame.paragraphs[5].text = '時雨量最高'+str(MD72DATA[2])+'mm('+MD72DATA[3]+'時)，累積降雨量'+str(MD72DATA[1])+'mm'
body_shape[2].text_frame.paragraphs[5].font.size = Pt(16)
body_shape[2].text_frame.paragraphs[5].font.color.rgb = RGBColor(255, 0, 0)
body_shape[2].text_frame.paragraphs[7].text = '時雨量最高'+str(AG72DATA[2])+'mm('+AG72DATA[3]+'時)，累積降雨量'+str(AG72DATA[1])+'mm'
body_shape[2].text_frame.paragraphs[7].font.size = Pt(16)
body_shape[2].text_frame.paragraphs[7].font.color.rgb = RGBColor(255, 0, 0)
body_shape[2].text_frame.paragraphs[9].text = '時雨量最高'+str(ZS72DATA[2])+'mm('+ZS72DATA[3]+'時)，累積降雨量'+str(ZS72DATA[1])+'mm'
body_shape[2].text_frame.paragraphs[9].font.size = Pt(16)
body_shape[2].text_frame.paragraphs[9].font.color.rgb = RGBColor(255, 0, 0)
body_shape[2].text_frame.paragraphs[11].text = '時雨量最高'+str(GP72DATA[2])+'mm('+GP72DATA[3]+'時)，累積降雨量'+str(GP72DATA[1])+'mm'
body_shape[2].text_frame.paragraphs[11].font.size = Pt(16)
body_shape[2].text_frame.paragraphs[11].font.color.rgb = RGBColor(255, 0, 0)
#-----第2頁文字----
#-----第3頁文字----
slide = prs.slides[2]
body_shape = slide.shapes
body_shape[1].text_frame.text = '南水局水庫集水區平均：一日預報(' + time.strftime('%m')+'月'+time.strftime('%d') \
+'日'+time.strftime('%H')+'時~'+time.strftime('%m')+'月'+str(int(time.strftime('%d'))+1)+'日'+str(int(time.strftime('%H'))-1)+'時)'

requestURL = 'http://localhost/QPESUMSQPF/'
print (requestURL)
js = urllib.request.urlopen(requestURL)
j_obj = json.loads(js.read())
for key in j_obj['values']:
  if key['code'] == 'Res0016':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    ZW24DATA = [name,sum(value),peak,peaktime]
  elif key['code'] == 'Res0022':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    MD24DATA = [name,sum(value),peak,peaktime]
  elif key['code'] == 'Res0020':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    AG24DATA = [name,sum(value),peak,peaktime]
  elif key['code'] == '1730w002':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    ZS24DATA = [name,sum(value),peak,peaktime]
  elif key['code'] == 'weir00002':
    name = key['code']
    value = [int(round(n)) for n in NulltoZero(key['value'])]
    if sum(value) != 0:
      peak = max(value)
      peaktime = (datetime.datetime.now() + datetime.timedelta(hours=int(reverseMax(value)))).strftime('%m/%d %H')
    else:
      peak = '0'
      peaktime = ''
    GP24DATA = [name,sum(value),peak,peaktime]
body_shape[7].text_frame.text = '時雨量最高'+str(MD24DATA[2])+'mm('+MD24DATA[3]+'時)\n累積降雨量'+str(MD24DATA[1])+'mm'
body_shape[8].text_frame.text = '時雨量最高'+str(AG24DATA[2])+'mm('+AG24DATA[3]+'時)\n累積降雨量'+str(AG24DATA[1])+'mm'
body_shape[9].text_frame.text = '時雨量最高'+str(ZW24DATA[2])+'mm('+ZW24DATA[3]+'時)\n累積降雨量'+str(ZW24DATA[1])+'mm'
#-----第3頁文字----
#-----第4頁文字----
slide = prs.slides[3]
body_shape = slide.shapes
body_shape[1].text_frame.text = '南水局水庫集水區平均：三日預報(' + time.strftime('%m')+'月'+time.strftime('%d') \
+'日'+time.strftime('%H')+'時~'+time.strftime('%m')+'月'+str(int(time.strftime('%d'))+3)+'日'+str(int(time.strftime('%H'))-1)+'時)'
body_shape[7].text_frame.text = '時雨量最高'+str(MD72DATA[2])+'mm('+MD72DATA[3]+'時)\n累積降雨量'+str(MD72DATA[1])+'mm'
body_shape[8].text_frame.text = '時雨量最高'+str(AG72DATA[2])+'mm('+AG72DATA[3]+'時)\n累積降雨量'+str(AG72DATA[1])+'mm'
body_shape[9].text_frame.text = '時雨量最高'+str(ZW72DATA[2])+'mm('+ZW72DATA[3]+'時)\n累積降雨量'+str(ZW72DATA[1])+'mm'
#-----第4頁文字----
#------簡報第3頁圖片--------
"""涉及前公司網頁爬蟲，不提供測試
browser.get('')
element =  wait.until(EC.presence_of_element_located((By.XPATH,"//div[@id='div-data-body']//div[@class='col-12 mb-3']")))
browser.save_screenshot(clpfigPath+'ZWQPF.png')
js="var q=document.documentElement.scrollTop=1000"
browser.execute_script(js)
browser.save_screenshot(clpfigPath+'MDQPF.png')
js="var q=document.documentElement.scrollTop=2000"
browser.execute_script(js)
browser.save_screenshot(clpfigPath+'AGQPF.png')
Imgclip=Image.open(clpfigPath+'ZWQPF.png')
ImgCrop = Imgclip.crop((150,550,1230,860))
ImgCrop.save(clpfigPath+'ZWQPFclip.png')
Imgclip=Image.open(clpfigPath+'MDQPF.png')
ImgCrop = Imgclip.crop((150,340,1230,650))
ImgCrop.save(clpfigPath+'MDQPFclip.png')
Imgclip=Image.open(clpfigPath+'AGQPF.png')
ImgCrop = Imgclip.crop((150,140,1230,450))
ImgCrop.save(clpfigPath+'AGQPFclip.png')
P4_1Path = clpfigPath+'ZWQPFclip.png'
P4_2Path = clpfigPath+'MDQPFclip.png'
P4_3Path = clpfigPath+'AGQPFclip.png'
#------簡報第3頁貼圖--------
slide = prs.slides[2]
shape = slide.shapes
Left = Cm(1.2)
Top1 = Cm(3.5)
Top2 = Cm(8.5)
Top3 = Cm(13.5)
height = Cm(4.9)
Pic1 = slide.shapes.add_picture(P4_1Path, Left, Top1, height=height)
slide.shapes._spTree.remove(Pic1._element)
slide.shapes._spTree.insert(2 , Pic1._element)
Pic2 = slide.shapes.add_picture(P4_2Path, Left, Top2, height=height)
slide.shapes._spTree.remove(Pic2._element)
slide.shapes._spTree.insert(2 , Pic2._element)
Pic3 = slide.shapes.add_picture(P4_3Path, Left, Top3, height=height)
slide.shapes._spTree.remove(Pic2._element)
slide.shapes._spTree.insert(2 , Pic2._element)
#------簡報第4頁圖片--------
select = Select(browser.find_element_by_id('select-fcstrainfall'))
time.sleep(1)
select.select_by_visible_text('QPESUMS_WRF')
browser.find_element_by_id("btn-find").click()
element =  wait.until(EC.text_to_be_present_in_element((By.XPATH,"//div[@id='div-data-body'][1][1]//div[@class='card-header text-center']"),'曾文水庫集水區QPESUMS_WRF降雨預報組體圖'))
browser.save_screenshot(clpfigPath+'ZWWRF.png')
js="var q=document.documentElement.scrollTop=1000"
browser.execute_script(js)
browser.save_screenshot(clpfigPath+'MDWRF.png')
js="var q=document.documentElement.scrollTop=2000"
browser.execute_script(js)
browser.save_screenshot(clpfigPath+'AGWRF.png')
Imgclip=Image.open(clpfigPath+'ZWWRF.png')
ImgCrop = Imgclip.crop((150,550,1230,860))
ImgCrop.save(clpfigPath+'ZWWRFclip.png')
Imgclip=Image.open(clpfigPath+'MDWRF.png')
ImgCrop = Imgclip.crop((150,340,1230,650))
ImgCrop.save(clpfigPath+'MDWRFclip.png')
Imgclip=Image.open(clpfigPath+'AGWRF.png')
ImgCrop = Imgclip.crop((150,140,1230,450))
ImgCrop.save(clpfigPath+'AGWRFclip.png')
P4_1Path = clpfigPath+'ZWWRFclip.png'
P4_2Path = clpfigPath+'MDWRFclip.png'
P4_3Path = clpfigPath+'AGWRFclip.png'
#------簡報第4頁貼圖--------
slide = prs.slides[3]
shape = slide.shapes
Left = Cm(1.2)
Top1 = Cm(3.5)
Top2 = Cm(8.5)
Top3 = Cm(13.5)
height = Cm(4.9)
Pic1 = slide.shapes.add_picture(P4_1Path, Left, Top1, height=height)
slide.shapes._spTree.remove(Pic1._element)
slide.shapes._spTree.insert(2 , Pic1._element)
Pic2 = slide.shapes.add_picture(P4_2Path, Left, Top2, height=height)
slide.shapes._spTree.remove(Pic2._element)
slide.shapes._spTree.insert(2 , Pic2._element)
Pic3 = slide.shapes.add_picture(P4_3Path, Left, Top3, height=height)
slide.shapes._spTree.remove(Pic2._element)
slide.shapes._spTree.insert(2 , Pic2._element)
"""
#------簡報第5頁圖片--------
r = requests.get('https://www.cwb.gov.tw/Data/fcst_img/QPF_ChFcstPrecip_12_12.png')
with open(clpfigPath+'QPF12_12.jpg','wb') as f:
    f.write(r.content)
r = requests.get('https://www.cwb.gov.tw/Data/fcst_img/QPF_ChFcstPrecip_12_24.png')
with open(clpfigPath+'QPF12_24.jpg','wb') as f:
    f.write(r.content)
#------簡報第5頁貼圖--------
P5_1Path = clpfigPath+'QPF12_12.jpg'
P5_2Path = clpfigPath+'QPF12_24.jpg'
slide = prs.slides[4]
shape = slide.shapes
Left1 = Cm(2.5)
Left2 = Cm(13)
Top = Cm(4)
height = Cm(12)
Pic1 = slide.shapes.add_picture(P5_1Path, Left1, Top, height=height)
slide.shapes._spTree.remove(Pic1._element)
slide.shapes._spTree.insert(2 , Pic1._element)
Pic2 = slide.shapes.add_picture(P5_2Path, Left2, Top, height=height)
slide.shapes._spTree.remove(Pic2._element)
slide.shapes._spTree.insert(2 , Pic2._element)
#-----簡報第7、8頁圖片----
if  int(datetime.datetime.now().strftime('%H')) >=8 and int(datetime.datetime.now().strftime('%H')) <14:
 timenow = datetime.datetime.now().replace(hour=8)
elif  int(datetime.datetime.now().strftime('%H')) >=14 and int(datetime.datetime.now().strftime('%H')) < 20:
 timenow = datetime.datetime.now().replace(hour=14)
elif  int(datetime.datetime.now().strftime('%H')) >= 20:
 timenow = datetime.datetime.now().replace(hour=20)
elif  int(datetime.datetime.now().strftime('%H')) < 2:
 timenow = datetime.datetime.now().replace(hour=20)
elif  int(datetime.datetime.now().strftime('%H')) >=2 and int(datetime.datetime.now().strftime('%H')) < 8:
 timenow = datetime.datetime.now().replace(hour=2)
dayminus20hour = (timenow - datetime.timedelta(hours=14)).strftime('%m%d%H')
dayminus20hourandyear = (timenow - datetime.timedelta(hours=14)).strftime('%Y%m%d%H')
dayminus1 = (datetime.datetime.now() - datetime.timedelta(days=1)).strftime('%m%d')
dayplus1 = (datetime.datetime.now() + datetime.timedelta(days=1)).strftime('%m%d')
dayplus2 = (datetime.datetime.now() + datetime.timedelta(days=2)).strftime('%m%d')
dayplus3 = (datetime.datetime.now() + datetime.timedelta(days=3)).strftime('%m%d')
dayplus4 = (datetime.datetime.now() + datetime.timedelta(days=4)).strftime('%m%d')
print (dayminus20hour)
r = requests.get('https://watch.ncdr.nat.gov.tw/00_Wxmap/5F1_NCDR_WRF_5km_precipitation/05days/'+time.strftime('%Y%m')+'/'+dayminus20hour+'_N01/rainfall/small-'+dayminus20hourandyear+'-'+time.strftime('%m%d')+'16-'+dayplus1+'16.gif')
with open(clpfigPath+'NCDR.jpg','wb') as f:
    f.write(r.content)
r = requests.get('https://watch.ncdr.nat.gov.tw/00_Wxmap/5F1_NCDR_WRF_5km_precipitation/05days/'+time.strftime('%Y%m')+'/'+dayminus20hour+'_N01/rainfall/small-'+dayminus20hourandyear+'-'+dayplus1+'16-'+dayplus2+'16.gif')
with open(clpfigPath+'NCDR2.jpg','wb') as f:
    f.write(r.content)
r = requests.get('https://watch.ncdr.nat.gov.tw/00_Wxmap/5F1_NCDR_WRF_5km_precipitation/05days/'+time.strftime('%Y%m')+'/'+dayminus20hour+'_N01/rainfall/small-'+dayminus20hourandyear+'-'+dayplus2+'16-'+dayplus3+'16.gif')
with open(clpfigPath+'NCDR3.jpg','wb') as f:
    f.write(r.content)
r = requests.get('https://watch.ncdr.nat.gov.tw/00_Wxmap/5F1_NCDR_WRF_5km_precipitation/05days/'+time.strftime('%Y%m')+'/'+dayminus20hour+'_N01/rainfall/small-'+dayminus20hourandyear+'-'+dayplus3+'16-'+dayplus4+'16.gif')
with open(clpfigPath+'NCDR4.jpg','wb') as f:
    f.write(r.content)
#-----簡報第7、8頁貼圖----
P7_1Path = clpfigPath+'NCDR.jpg'
P7_2Path = clpfigPath+'NCDR2.jpg'
P8_1Path = clpfigPath+'NCDR3.jpg'
P8_2Path = clpfigPath+'NCDR4.jpg'
slide = prs.slides[6]
shape = slide.shapes
Left1 = Cm(2)
Left2 = Cm(13.8)
Top = Cm(4)
height = Cm(13)
Pic1 = slide.shapes.add_picture(P7_1Path, Left1, Top, height=height)
slide.shapes._spTree.remove(Pic1._element)
slide.shapes._spTree.insert(2 , Pic1._element)
Pic2 = slide.shapes.add_picture(P7_2Path, Left2, Top, height=height)
slide.shapes._spTree.remove(Pic2._element)
slide.shapes._spTree.insert(2 , Pic2._element)
slide = prs.slides[7]
Pic3 = slide.shapes.add_picture(P8_1Path, Left1, Top, height=height)
slide.shapes._spTree.remove(Pic3._element)
slide.shapes._spTree.insert(2 , Pic3._element)
Pic4 = slide.shapes.add_picture(P8_2Path, Left2, Top, height=height)
slide.shapes._spTree.remove(Pic4._element)
slide.shapes._spTree.insert(2 , Pic4._element)
#-----第6頁文字----
slide = prs.slides[5]
body_shape = slide.shapes
body_shape[1].text_frame.paragraphs[0].text = (datetime.datetime.now() + datetime.timedelta(days=1)).strftime('%m/%d日')+'\v曾文水庫累積降雨最高100mm、牡丹水庫200mm、阿公店水庫30mm。'
body_shape[1].text_frame.paragraphs[0].font.size = Pt(22)
body_shape[1].text_frame.paragraphs[1].text = (datetime.datetime.now() + datetime.timedelta(days=2)).strftime('%m/%d日')+'\v曾文水庫累積降雨最高80mm、牡丹水庫150mm、阿公店水庫40mm。'
body_shape[1].text_frame.paragraphs[1].font.size = Pt(22)
body_shape[1].text_frame.paragraphs[2].text = (datetime.datetime.now() + datetime.timedelta(days=3)).strftime('%m/%d日')+'\v曾文水庫累積降雨最高170mm、牡丹水庫70mm、阿公店水庫50mm。'
body_shape[1].text_frame.paragraphs[2].font.size = Pt(22)
body_shape[1].text_frame.paragraphs[3].text = (datetime.datetime.now() + datetime.timedelta(days=4)).strftime('%m/%d日')+'\v曾文水庫累積降雨最高110mm、牡丹水庫80mm、阿公店水庫60mm。'
body_shape[1].text_frame.paragraphs[3].font.size = Pt(22)
#-----第6頁文字----
prs.save(time.strftime('%Y%m%d%H')+u'水情研判簡報.pptx')
browser.quit()
print ('pptOutput succeed')